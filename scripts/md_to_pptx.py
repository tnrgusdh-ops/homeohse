#!/usr/bin/env python3
"""Convert a simple Markdown slides file into PPTX.

Slides are separated by a line with only three dashes: ---
The first heading in each slide (a line starting with # or ##) is used as the slide title.
Images in Markdown syntax ![alt](path) will be added (if file exists).

Usage: python scripts/md_to_pptx.py input.md output.pptx [--template template.pptx]
"""
import re
import sys
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception as e:
    print("Missing dependency 'python-pptx'. Install with: pip install python-pptx", file=sys.stderr)
    raise

IMG_RE = re.compile(r"!\[(.*?)\]\((.*?)\)")


def split_slides(md_text: str):
    parts = re.split(r"^---$", md_text, flags=re.MULTILINE)
    return [p.strip() for p in parts if p.strip()]


def extract_title_and_body(slide_md: str):
    lines = [l for l in slide_md.splitlines()]
    title = ""
    body_lines = []
    for i, line in enumerate(lines):
        if re.match(r"^#{1,6}\s+", line):
            title = re.sub(r"^#{1,6}\s+", "", line).strip()
            body_lines = [l for l in lines[i+1:] if l.strip()]
            break
    else:
        # no heading found, use first non-empty line as title
        for i, line in enumerate(lines):
            if line.strip():
                title = line.strip()
                body_lines = [l for l in lines[i+1:] if l.strip()]
                break
    return title, "\n".join(body_lines)


def add_slide(prs: Presentation, title: str, body: str, base_path: Path):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    if title:
        slide.shapes.title.text = title
    if body:
        # handle images in body separately and remove image lines from text
        images = []
        def _img_repl(m):
            images.append((m.group(1), m.group(2)))
            return ""  # remove image tag from text
        text_without_imgs = IMG_RE.sub(_img_repl, body).strip()

        if text_without_imgs:
            try:
                tx_box = slide.placeholders[1].text_frame
                tx_box.text = text_without_imgs
            except Exception:
                pass

        # add images below the text
        left = Inches(1)
        top = Inches(2.0)
        max_h = Inches(4)
        for alt, relpath in images:
            img_path = Path(relpath)
            if not img_path.is_absolute():
                img_path = (base_path / relpath).resolve()
            if img_path.exists():
                try:
                    pic = slide.shapes.add_picture(str(img_path), left, top, height=max_h)
                    top = top + pic.height + Inches(0.2)
                except Exception as e:
                    print(f"Warning: failed to add image {img_path}: {e}", file=sys.stderr)
            else:
                print(f"Warning: image not found: {img_path}", file=sys.stderr)


def main():
    p = argparse.ArgumentParser(description="Convert simple Markdown slides to PPTX")
    p.add_argument("input_md", help="Input Markdown file")
    p.add_argument("output_pptx", help="Output PPTX file")
    p.add_argument("--template", help="Optional PPTX template to base slides on")
    args = p.parse_args()

    in_path = Path(args.input_md)
    if not in_path.exists():
        print(f"Input file not found: {in_path}", file=sys.stderr)
        sys.exit(2)

    base_path = in_path.parent

    text = in_path.read_text(encoding="utf-8")
    slides_md = split_slides(text)

    if args.template:
        tpl = Path(args.template)
        if not tpl.exists():
            print(f"Template not found: {tpl}", file=sys.stderr)
            sys.exit(2)
        prs = Presentation(str(tpl))
    else:
        prs = Presentation()

    for s in slides_md:
        title, body = extract_title_and_body(s)
        add_slide(prs, title, body, base_path)

    prs.save(args.output_pptx)
    print(f"Saved: {args.output_pptx}")


if __name__ == "__main__":
    main()
